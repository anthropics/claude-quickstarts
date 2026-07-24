#!/usr/bin/env python3
"""Assemble the deal-room KG cookbook — REAL-CORPUS customer preview edition."""
import json
from pathlib import Path

cells = []


def md(src):
    cells.append({"cell_type": "markdown", "metadata": {}, "source": src.splitlines(keepends=True)})


def code(src):
    cells.append({"cell_type": "code", "execution_count": None, "metadata": {},
                  "outputs": [], "source": src.splitlines(keepends=True)})


md("""# Build a knowledge wiki that cuts agentic retrieval costs

When agents answer questions over a document corpus by searching it, every question pays the same tax: re-read, re-search, re-derive. The costs compound in two ways — tokens (each answer re-reads material a previous answer already digested) and correctness (the hard facts hide in the seams between documents: a number that only exists in a chart, two metrics with the same name on different bases, one entity behind three different names).

Extracting the corpus **once** into a knowledge wiki, then letting agents answer from the wiki, fixes both. The pattern fits any domain where agents repeatedly query a slow-changing document set — for example:

- **M&A / due diligence** — a deal team asking hundreds of questions against a data room (the worked example below)
- **Legal discovery** — one matter's corpus serving months of queries from the whole team
- **Support over product docs** — agents answering from manuals, changelogs, and past tickets

**"Isn't this just document indexing / RAG?"** No — and the difference is
what this notebook demonstrates. An index (vector, keyword, or graph DB)
helps an agent *find passages of the original documents*; the agent still
re-reads and re-synthesizes them on every query, and answers that live in
the *seams between documents* — an entity behind three different names, a
figure that changed mid-process, two metrics with the same label on
different bases — don't exist in any retrievable passage at all. Here,
Claude reads the corpus **once** and writes down its *conclusions*, in the
[Agent Knowledge format](https://github.com/anthropics/skills) - markdown
nodes with YAML frontmatter, [[wikilinks]] as edges (the same thin standard
behind Skills):
entities resolved across aliases, facts temporally ordered with
supersession, conflicts flagged, every claim carrying provenance. Queries
then read conclusions instead of re-deriving them. Two practical bonuses:
the store is plain markdown in a managed memory store — inspectable and
editable by humans, no embedding infrastructure to run — and the baseline
we compare against below is *stronger* than typical RAG: a full agent with
unlimited search over the complete corpus.

In our full-scale ablation of this pattern, the wiki answered a 21-question judged battery at **96%** quality versus **86%** for that agentic-search baseline — at **one fifth** of the per-question token cost.

**The worked example** is the first scenario, chosen because its documents are public and maximally adversarial: 26 real SEC filings from the 2024 take-private of Squarespace, Inc. by Permira (fetched live from EDGAR), a real financial-advisor board deck filed as scanned slide images, plus two clearly-marked synthetic files standing in for the internal documents a real deal room adds. It carries every failure mode worth testing — a deal price that changed mid-process, look-alike shell entities, code names, and chart-only numbers. Nothing in the pipeline is deal-specific: swap the corpus and the schema vocabulary, and the same six steps apply.

**By the end of this cookbook, you'll be able to:**

- Assemble a real, reproducible document room from public sources — at three cost tiers (26 / 37 / 42 documents)
- Normalize PDFs (charts included, via native document blocks) and Word files into a provenance-stamped corpus
- Build a knowledge wiki in a memory store with extraction agents — the schema travels on the store attachment
- Resolve the wiki's open questions *before* consolidation, then run one **steered dream** to organize everything
- Query the wiki read-only at a fraction of raw-document cost, with a citation on every fact

""")

md("""## The architecture

Two flows share exactly one thing — the knowledge wiki store. The build flow runs once per deal: extraction (with gap resolution) writes the wiki, and a steered dream consolidates it. The query flow runs per question: analysts attach the consolidated store read-only. The dashed arrow is the learning loop you'll run in Step 5: analyst session transcripts feed a periodic usage dream, so the wiki reorganizes around how it's actually queried — and analyst misses become open questions for the next build cycle.

<img alt="Architecture: build flow (extract and resolve, then dream) and query flow sharing one knowledge wiki store" src="data:image/svg+xml;base64,PHN2ZyB4bWxucz0iaHR0cDovL3d3dy53My5vcmcvMjAwMC9zdmciIHZpZXdCb3g9IjAgMCA5NjAgNzAwIiBmb250LWZhbWlseT0iR2VvcmdpYSwgc2VyaWYiPgo8c3R5bGU+OnJvb3R7LS1hY2M6IzJhNzhkNjstLWJnOiNmY2ZjZmI7LS1jYXJkOiNmZmZmZmY7LS1pbms6IzBiMGIwYjstLWxpbmU6I2RlZGNkNTstLW11dGVkOiM3MzcyNmM7LS1zMTojMmE3OGQ2Oy0tczM6I2VkYTEwMDstLXNlYzojNTI1MTRlO30KLmQyLWNhbnZhc3tmaWxsOmNvbG9yLW1peChpbiBzcmdiLHZhcigtLXMzKSA2JSx2YXIoLS1iZykpO30uZDItYm94e2ZpbGw6dmFyKC0tY2FyZCk7c3Ryb2tlOnZhcigtLWluayk7c3Ryb2tlLXdpZHRoOjEuMzt9LmQyLWhlcm97ZmlsbDojZWI2ODM0O3N0cm9rZTp2YXIoLS1pbmspO3N0cm9rZS13aWR0aDoxLjY7c3Ryb2tlLWRhc2hhcnJheTo2IDQ7fS5kMi1idWJibGV7ZmlsbDpjb2xvci1taXgoaW4gc3JnYix2YXIoLS1zMSkgMTQlLHZhcigtLWNhcmQpKTtzdHJva2U6Y29sb3ItbWl4KGluIHNyZ2IsdmFyKC0tczEpIDQ1JSx2YXIoLS1saW5lKSk7c3Ryb2tlLXdpZHRoOjEuMjt9LmQyLXdlYntmaWxsOmNvbG9yLW1peChpbiBzcmdiLHZhcigtLXMzKSAxNCUsdmFyKC0tY2FyZCkpO3N0cm9rZTp2YXIoLS1saW5lKTtzdHJva2Utd2lkdGg6MS4yO30uZDItdHtmb250OjcwMCAxN3B4IENoYXJ0ZXIsR2VvcmdpYSxzZXJpZjtmaWxsOnZhcigtLWluayk7fS5kMi1oZXJvLXR7Zm9udDo3MDAgMTlweCBDaGFydGVyLEdlb3JnaWEsc2VyaWY7ZmlsbDojZmZmO30uZDItbXtmb250Oml0YWxpYyAxMXB4IENoYXJ0ZXIsR2VvcmdpYSxzZXJpZjtmaWxsOnZhcigtLW11dGVkKTt9LmQyLWhlcm8tbXtmb250Oml0YWxpYyAxMXB4IENoYXJ0ZXIsR2VvcmdpYSxzZXJpZjtmaWxsOiNmZmUzZDY7fS5kMi1ie2ZvbnQ6MTJweCBDaGFydGVyLEdlb3JnaWEsc2VyaWY7ZmlsbDp2YXIoLS1zZWMpO30uZDItaGVyby1ie2ZvbnQ6MTJweCBDaGFydGVyLEdlb3JnaWEsc2VyaWY7ZmlsbDojZmZmO30uZDItbGJse2ZvbnQ6aXRhbGljIDExLjVweCBDaGFydGVyLEdlb3JnaWEsc2VyaWY7ZmlsbDp2YXIoLS1tdXRlZCk7fS5kMi1jYXB7Zm9udDppdGFsaWMgMTJweCBDaGFydGVyLEdlb3JnaWEsc2VyaWY7ZmlsbDp2YXIoLS1tdXRlZCk7fS5kMi1mbHtzdHJva2U6dmFyKC0taW5rKTtzdHJva2Utd2lkdGg6MS42O30uZDItZmwye3N0cm9rZTp2YXIoLS1tdXRlZCk7c3Ryb2tlLXdpZHRoOjEuNDtvcGFjaXR5Oi43NTt9LmQyLWRie2ZpbGw6dmFyKC0tY2FyZCk7c3Ryb2tlOnZhcigtLWluayk7c3Ryb2tlLXdpZHRoOjEuMzt9LmQyLWRiLWtne2ZpbGw6Y29sb3ItbWl4KGluIHNyZ2IsdmFyKC0tczEpIDEwJSx2YXIoLS1jYXJkKSk7c3Ryb2tlOnZhcigtLWFjYyk7c3Ryb2tlLXdpZHRoOjI7fS5kMi1iYW5ke2ZpbGw6bm9uZTtzdHJva2U6dmFyKC0tbXV0ZWQpO3N0cm9rZS13aWR0aDoxO3N0cm9rZS1kYXNoYXJyYXk6MyA1O29wYWNpdHk6LjY7fS5kMi1iYW5kbGJse2ZvbnQ6NzAwIDEycHggdWktbW9ub3NwYWNlLE1lbmxvLG1vbm9zcGFjZTtsZXR0ZXItc3BhY2luZzouMDhlbTtmaWxsOnZhcigtLW11dGVkKTt9LmQyLXN0ZXB7ZmlsbDp2YXIoLS1pbmspO30uZDItc3RlcG57Zm9udDo3MDAgMTJweCB1aS1tb25vc3BhY2UsTWVubG8sbW9ub3NwYWNlO2ZpbGw6dmFyKC0tYmcpO30Kc3Zne2JhY2tncm91bmQ6I2ZhZjdmMn08L3N0eWxlPgoKPGRlZnM+CjxtYXJrZXIgaWQ9ImFoIiB2aWV3Qm94PSIwIDAgMTAgMTAiIHJlZlg9IjkiIHJlZlk9IjUiIG1hcmtlcldpZHRoPSI4IiBtYXJrZXJIZWlnaHQ9IjgiIG9yaWVudD0iYXV0by1zdGFydC1yZXZlcnNlIj48cGF0aCBkPSJNMCwwTDEwLDVMMCwxMHoiIGZpbGw9InZhcigtLWluaykiLz48L21hcmtlcj4KPG1hcmtlciBpZD0iYWgyIiB2aWV3Qm94PSIwIDAgMTAgMTAiIHJlZlg9IjkiIHJlZlk9IjUiIG1hcmtlcldpZHRoPSI4IiBtYXJrZXJIZWlnaHQ9IjgiIG9yaWVudD0iYXV0by1zdGFydC1yZXZlcnNlIj48cGF0aCBkPSJNMCwwTDEwLDVMMCwxMHoiIGZpbGw9InZhcigtLW11dGVkKSIvPjwvbWFya2VyPgo8L2RlZnM+CjxyZWN0IGNsYXNzPSJkMi1jYW52YXMiIHg9IjAiIHk9IjAiIHdpZHRoPSI5NjAiIGhlaWdodD0iNzAwIiByeD0iMTIiLz4KCjxyZWN0IGNsYXNzPSJkMi1iYW5kIiB4PSIxNiIgeT0iMTYiIHdpZHRoPSI5MjgiIGhlaWdodD0iMjU2IiByeD0iMTAiLz4KPHRleHQgY2xhc3M9ImQyLWJhbmRsYmwiIHg9IjMyIiB5PSI0MCI+RkxPVyAxICZtaWRkb3Q7IEJVSUxEIFRIRSBHUkFQSCAmbWRhc2g7IGV4dHJhY3QgJmFtcDsgcmVzb2x2ZSwgdGhlbiBkcmVhbTwvdGV4dD4KCjxyZWN0IGNsYXNzPSJkMi1iYW5kIiB4PSIxNiIgeT0iNTE2IiB3aWR0aD0iOTI4IiBoZWlnaHQ9IjE2OCIgcng9IjEwIi8+Cjx0ZXh0IGNsYXNzPSJkMi1iYW5kbGJsIiB4PSIzMiIgeT0iNTQwIj5GTE9XIDIgJm1pZGRvdDsgUVVFUlkgVEhFIEdSQVBIICZtZGFzaDsgZXZlcnkgYW5hbHlzdCBxdWVzdGlvbjwvdGV4dD4KCjxwYXRoIGNsYXNzPSJkMi1kYiIgZD0iTSA1MCw5NiB2IDk2IGEgNzUuMCwxMiAwIDAgMCAxNTAsMCB2IC05NiIvPjxlbGxpcHNlIGNsYXNzPSJkMi1kYiIgY3g9IjEyNS4wIiBjeT0iOTYiIHJ4PSI3NS4wIiByeT0iMTIiLz4KPHRleHQgY2xhc3M9ImQyLXQiIHg9IjEyNSIgeT0iMTIyIiB0ZXh0LWFuY2hvcj0ibWlkZGxlIiBzdHlsZT0iZm9udC1zaXplOjE0cHgiPkRlYWwgZGF0YSByb29tPC90ZXh0Pgo8dGV4dCBjbGFzcz0iZDItYiIgeD0iMTI1IiB5PSIxNDIiIHRleHQtYW5jaG9yPSJtaWRkbGUiPmZpbGluZ3MgwrcgZGVjayDCtyBub3RlczwvdGV4dD4KPHRleHQgY2xhc3M9ImQyLW0iIHg9IjEyNSIgeT0iMTYwIiB0ZXh0LWFuY2hvcj0ibWlkZGxlIj5mZXRjaGVkIGZyb20gcHVibGljIHNvdXJjZXM8L3RleHQ+Cgo8cmVjdCBjbGFzcz0iZDItYm94IiB4PSIyNzAiIHk9IjY2IiB3aWR0aD0iMjQwIiBoZWlnaHQ9IjE1NiIgcng9IjE0Ii8+CjxjaXJjbGUgY2xhc3M9ImQyLXN0ZXAiIGN4PSIyNzAiIGN5PSI2NiIgcj0iMTEiLz48dGV4dCBjbGFzcz0iZDItc3RlcG4iIHg9IjI3MCIgeT0iNzAiIHRleHQtYW5jaG9yPSJtaWRkbGUiPjE8L3RleHQ+Cjx0ZXh0IGNsYXNzPSJkMi10IiB4PSIzOTAiIHk9Ijk2IiB0ZXh0LWFuY2hvcj0ibWlkZGxlIj5FeHRyYWN0IGFuZCBSZXNvbHZlIEFnZW50PC90ZXh0Pgo8dGV4dCBjbGFzcz0iZDItbSIgeD0iMzkwIiB5PSIxMTQiIHRleHQtYW5jaG9yPSJtaWRkbGUiPlNvbm5ldCA1ICZtaWRkb3Q7IG9uZSBzZXNzaW9uIHBlciBkb2MsPC90ZXh0Pgo8dGV4dCBjbGFzcz0iZDItbSIgeD0iMzkwIiB5PSIxMzAiIHRleHQtYW5jaG9yPSJtaWRkbGUiPnRoZW4gYSByZXNvbHZlIHBhc3Mgb3ZlciB0aGUgZ2FwIGxpc3Q8L3RleHQ+Cjx0ZXh0IGNsYXNzPSJkMi1iIiB4PSIzOTAiIHk9IjE1NCIgdGV4dC1hbmNob3I9Im1pZGRsZSI+ZmlsZXMgZmFjdHMgd2l0aCBbc291cmNlIHwgYXMtb2ZdPC90ZXh0Pgo8dGV4dCBjbGFzcz0iZDItYiIgeD0iMzkwIiB5PSIxNzIiIHRleHQtYW5jaG9yPSJtaWRkbGUiPmdhcHMgJnJhcnI7IG9wZW5fcXVlc3Rpb25zLm1kLCB0aGVuPC90ZXh0Pgo8dGV4dCBjbGFzcz0iZDItYiIgeD0iMzkwIiB5PSIxOTAiIHRleHQtYW5jaG9yPSJtaWRkbGUiPnJlc2VhcmNoZWQgb24gU0VDIEVER0FSICh3ZWIgdG9vbHMpPC90ZXh0PgoKPHJlY3QgY2xhc3M9ImQyLWhlcm8iIHg9IjcwMCIgeT0iNjYiIHdpZHRoPSIyMzAiIGhlaWdodD0iMTMyIiByeD0iMTQiLz4KPGNpcmNsZSBjeD0iNzAwIiBjeT0iNjYiIHI9IjExIiBmaWxsPSIjZmZmIiBzdHJva2U9InZhcigtLWluaykiIHN0cm9rZS13aWR0aD0iMSIvPjx0ZXh0IHg9IjcwMCIgeT0iNzAiIHRleHQtYW5jaG9yPSJtaWRkbGUiIHN0eWxlPSJmb250OjcwMCAxMnB4IHVpLW1vbm9zcGFjZSxNZW5sbyxtb25vc3BhY2U7ZmlsbDojZWI2ODM0Ij4yPC90ZXh0Pgo8dGV4dCBjbGFzcz0iZDItaGVyby10IiB4PSI4MTUiIHk9IjEwMCIgdGV4dC1hbmNob3I9Im1pZGRsZSI+U3RlZXJlZCBkcmVhbTwvdGV4dD4KPHRleHQgY2xhc3M9ImQyLWhlcm8tbSIgeD0iODE1IiB5PSIxMTgiIHRleHQtYW5jaG9yPSJtaWRkbGUiPlBPU1QgL3YxL2RyZWFtcyAmbWlkZG90OyBjbGF1ZGUtc29ubmV0LTU8L3RleHQ+Cjx0ZXh0IGNsYXNzPSJkMi1oZXJvLWIiIHg9IjgxNSIgeT0iMTQyIiB0ZXh0LWFuY2hvcj0ibWlkZGxlIj5kZWR1cGUgJm1pZGRvdDsgaW5kZXggJm1pZGRvdDsgcGxhY2UgZmFjdHM8L3RleHQ+Cjx0ZXh0IGNsYXNzPSJkMi1oZXJvLWIiIHg9IjgxNSIgeT0iMTYwIiB0ZXh0LWFuY2hvcj0ibWlkZGxlIj53cml0ZSByYW5rZWQgZXNjYWxhdGlvbnMubWQ8L3RleHQ+Cjx0ZXh0IGNsYXNzPSJkMi1oZXJvLW0iIHg9IjgxNSIgeT0iMTgyIiB0ZXh0LWFuY2hvcj0ibWlkZGxlIj5yZWFkcyB0aGUgZ3JhcGggKyB0cmFuc2NyaXB0czwvdGV4dD4KCjxwYXRoIGNsYXNzPSJkMi1kYi1rZyIgZD0iTSA0MTAsMzYwIHYgMTE2IGEgOTUuMCwxMiAwIDAgMCAxOTAsMCB2IC0xMTYiLz48ZWxsaXBzZSBjbGFzcz0iZDItZGIta2ciIGN4PSI1MDUuMCIgY3k9IjM2MCIgcng9Ijk1LjAiIHJ5PSIxMiIvPgo8dGV4dCBjbGFzcz0iZDItdCIgeD0iNTA1IiB5PSIzOTYiIHRleHQtYW5jaG9yPSJtaWRkbGUiPktub3dsZWRnZSBHcmFwaDwvdGV4dD4KPHRleHQgY2xhc3M9ImQyLW0iIHg9IjUwNSIgeT0iNDI2IiB0ZXh0LWFuY2hvcj0ibWlkZGxlIj52ZXJzaW9uZWQgbWVtb3J5IHN0b3JlPC90ZXh0Pgo8dGV4dCBjbGFzcz0iZDItYiIgeD0iNTA1IiB5PSI0NDYiIHRleHQtYW5jaG9yPSJtaWRkbGUiPmluZGV4ICZtaWRkb3Q7IGVudGl0aWVzICZtaWRkb3Q7IGVzY2FsYXRpb25zPC90ZXh0PgoKPHJlY3QgY2xhc3M9ImQyLWJveCIgeD0iMzgwIiB5PSI1NjQiIHdpZHRoPSIyNTAiIGhlaWdodD0iMTA0IiByeD0iMTQiLz4KPGNpcmNsZSBjbGFzcz0iZDItc3RlcCIgY3g9IjM4MCIgY3k9IjU2NCIgcj0iMTEiLz48dGV4dCBjbGFzcz0iZDItc3RlcG4iIHg9IjM4MCIgeT0iNTY4IiB0ZXh0LWFuY2hvcj0ibWlkZGxlIj4zPC90ZXh0Pgo8dGV4dCBjbGFzcz0iZDItdCIgeD0iNTA1IiB5PSI1OTQiIHRleHQtYW5jaG9yPSJtaWRkbGUiPkFuYWx5c3QgYWdlbnRzPC90ZXh0Pgo8dGV4dCBjbGFzcz0iZDItbSIgeD0iNTA1IiB5PSI2MTIiIHRleHQtYW5jaG9yPSJtaWRkbGUiPmZyZXNoIENNQSBzZXNzaW9uIHBlciBxdWVzdGlvbjwvdGV4dD4KPHRleHQgY2xhc3M9ImQyLWIiIHg9IjUwNSIgeT0iNjM0IiB0ZXh0LWFuY2hvcj0ibWlkZGxlIj5hbnN3ZXIgZnJvbSB0aGUgZ3JhcGggb25seSw8L3RleHQ+Cjx0ZXh0IGNsYXNzPSJkMi1iIiB4PSI1MDUiIHk9IjY1MiIgdGV4dC1hbmNob3I9Im1pZGRsZSI+Y2l0ZSBbc291cmNlIHwgYXMtb2ZdIHByb3ZlbmFuY2U8L3RleHQ+Cgo8cGF0aCBjbGFzcz0iZDItYnViYmxlIiBkPSJNNjAsNTcyIGgyMjAgYTEwLDEwIDAgMCAxIDEwLDEwIHY1NiBhMTAsMTAgMCAwIDEgLTEwLDEwIGgtMTUwIGwtMTQsMTYgLTIsLTE2IGgtNTQgYTEwLDEwIDAgMCAxIC0xMCwtMTAgdi01NiBhMTAsMTAgMCAwIDEgMTAsLTEwIHoiLz4KPHRleHQgY2xhc3M9ImQyLW0iIHg9IjE4MCIgeT0iNTk2IiB0ZXh0LWFuY2hvcj0ibWlkZGxlIj5kZWFsLXRlYW0gcXVlc3Rpb25zPC90ZXh0Pgo8dGV4dCBjbGFzcz0iZDItYiIgeD0iMTgwIiB5PSI2MTgiIHRleHQtYW5jaG9yPSJtaWRkbGUiPiJ3aGljaCBFQklUREEgaXMgY29tcGFyYWJsZT8iPC90ZXh0Pgo8dGV4dCBjbGFzcz0iZDItYiIgeD0iMTgwIiB5PSI2MzYiIHRleHQtYW5jaG9yPSJtaWRkbGUiPiJhbnl0aGluZyB3b3J0aCBlc2NhbGF0aW5nPyI8L3RleHQ+Cgo8cmVjdCBjbGFzcz0iZDItd2ViIiB4PSI3MTAiIHk9IjU3MiIgd2lkdGg9IjIwMCIgaGVpZ2h0PSI4NCIgcng9IjE0Ii8+Cjx0ZXh0IGNsYXNzPSJkMi10IiB4PSI4MTAiIHk9IjYwMiIgdGV4dC1hbmNob3I9Im1pZGRsZSIgc3R5bGU9ImZvbnQtc2l6ZToxNHB4Ij5Tb3VyY2VkIGFuc3dlcnM8L3RleHQ+Cjx0ZXh0IGNsYXNzPSJkMi1iIiB4PSI4MTAiIHk9IjYyNCIgdGV4dC1hbmNob3I9Im1pZGRsZSI+ZXZlcnkgZmFjdCBjYXJyaWVzIHByb3ZlbmFuY2U7PC90ZXh0Pgo8dGV4dCBjbGFzcz0iZDItYiIgeD0iODEwIiB5PSI2NDIiIHRleHQtYW5jaG9yPSJtaWRkbGUiPiJub3QgaW4gdGhlIGRhdGEgcm9vbSIgb24gYSBtaXNzPC90ZXh0PgoKPGxpbmUgY2xhc3M9ImQyLWZsIiB4MT0iMjAwIiB5MT0iMTQ0IiB4Mj0iMjYyIiB5Mj0iMTQ0IiBtYXJrZXItZW5kPSJ1cmwoI2FoKSIvPgo8dGV4dCBjbGFzcz0iZDItbGJsIiB4PSIyMzEiIHk9IjEzMiIgdGV4dC1hbmNob3I9Im1pZGRsZSI+cmVhZCBvbmNlPC90ZXh0PgoKPGxpbmUgY2xhc3M9ImQyLWZsMiIgeDE9IjUxMCIgeTE9IjEzMiIgeDI9IjY5MiIgeTI9IjEzMiIgbWFya2VyLWVuZD0idXJsKCNhaDIpIi8+Cjx0ZXh0IGNsYXNzPSJkMi1sYmwiIHg9IjYwMSIgeT0iMTIwIiB0ZXh0LWFuY2hvcj0ibWlkZGxlIj50cmFuc2NyaXB0cyBmZWVkIHRoZSBkcmVhbTwvdGV4dD4KCjxsaW5lIGNsYXNzPSJkMi1mbCIgeDE9IjQyMCIgeTE9IjIyMiIgeDI9IjQ4MiIgeTI9IjM0MiIgbWFya2VyLWVuZD0idXJsKCNhaCkiLz4KPHRleHQgY2xhc3M9ImQyLWxibCIgeD0iNDAwIiB5PSIyOTIiIHRleHQtYW5jaG9yPSJlbmQiPndyaXRlcyBncmFwaCArIGdhcCByZXNvbHV0aW9uczwvdGV4dD4KCjxsaW5lIGNsYXNzPSJkMi1mbDIiIHgxPSI1NjAiIHkxPSIzNDQiIHgyPSI3NDIiIHkyPSIyMDIiIG1hcmtlci1lbmQ9InVybCgjYWgyKSIvPgo8dGV4dCBjbGFzcz0iZDItbGJsIiB4PSI2NzIiIHk9IjI5MiIgdGV4dC1hbmNob3I9InN0YXJ0Ij5yZWFkczwvdGV4dD4KCjxsaW5lIGNsYXNzPSJkMi1mbCIgeDE9Ijg1MCIgeTE9IjE5OCIgeDI9IjU5MiIgeTI9IjM1OCIgbWFya2VyLWVuZD0idXJsKCNhaCkiLz4KPHRleHQgY2xhc3M9ImQyLWxibCIgeD0iNzkwIiB5PSIyOTIiIHRleHQtYW5jaG9yPSJzdGFydCI+d3JpdGVzIGEgbmV3LDwvdGV4dD4KPHRleHQgY2xhc3M9ImQyLWxibCIgeD0iNzkwIiB5PSIzMDgiIHRleHQtYW5jaG9yPSJzdGFydCI+Y29uc29saWRhdGVkIHZlcnNpb248L3RleHQ+Cgo8bGluZSBjbGFzcz0iZDItZmwiIHgxPSI0OTAiIHkxPSI1MDAiIHgyPSI0OTAiIHkyPSI1NTYiIG1hcmtlci1lbmQ9InVybCgjYWgpIi8+Cjx0ZXh0IGNsYXNzPSJkMi1sYmwiIHg9IjQ3OCIgeT0iNTE2IiB0ZXh0LWFuY2hvcj0iZW5kIj5yZWFkX29ubHkgYXR0YWNoPC90ZXh0Pgo8bGluZSBjbGFzcz0iZDItZmwyIiB4MT0iNTQwIiB5MT0iNTU2IiB4Mj0iNTQwIiB5Mj0iNTAwIiBtYXJrZXItZW5kPSJ1cmwoI2FoMikiIHN0cm9rZS1kYXNoYXJyYXk9IjUgNCIvPgo8Y2lyY2xlIGNsYXNzPSJkMi1zdGVwIiBjeD0iNTQwIiBjeT0iNTUyIiByPSIxMSIvPjx0ZXh0IGNsYXNzPSJkMi1zdGVwbiIgeD0iNTQwIiB5PSI1NTYiIHRleHQtYW5jaG9yPSJtaWRkbGUiPjQ8L3RleHQ+Cjx0ZXh0IGNsYXNzPSJkMi1sYmwiIHg9IjU1NiIgeT0iNTI0Ij5kcmVhbSBvdmVyIGFuYWx5c3QgdHJhbnNjcmlwdHM8L3RleHQ+Cjx0ZXh0IGNsYXNzPSJkMi1sYmwiIHg9IjU1NiIgeT0iNTQwIj50aGUgZ3JhcGggbGVhcm5zIGZyb20gdXNhZ2U8L3RleHQ+Cgo8bGluZSBjbGFzcz0iZDItZmwiIHgxPSIyOTAiIHkxPSI2MTAiIHgyPSIzNzIiIHkyPSI2MTAiIG1hcmtlci1lbmQ9InVybCgjYWgpIi8+CjxsaW5lIGNsYXNzPSJkMi1mbCIgeDE9IjYzMCIgeTE9IjYxNCIgeDI9IjcwMiIgeTI9IjYxNCIgbWFya2VyLWVuZD0idXJsKCNhaCkiLz4KCjx0ZXh0IGNsYXNzPSJkMi1jYXAiIHg9IjQ4MCIgeT0iNjkyIiB0ZXh0LWFuY2hvcj0ibWlkZGxlIj50aGUgdHdvIGZsb3dzIHNoYXJlIGV4YWN0bHkgb25lIHRoaW5nOiB0aGUga25vd2xlZGdlIGdyYXBoIHN0b3JlLiBCdWlsZCBydW5zIG9uY2UgcGVyIGRlYWwgKCQ4NSBtZWFzdXJlZCk7IHF1ZXJ5IHJ1bnMgcGVyIHF1ZXN0aW9uICh+NjlrIHRva2VucykuPC90ZXh0Pgo8L3N2Zz4=" style="max-width:100%">

| Step | Flow | What happens |
|---|---|---|
| 1 | Build | Build the knowledge wiki from the source documents, and resolve any emergent open questions with supplemental research. |
| 2 | Build | Improve the knowledge wiki with async processing — the steered dream dedupes entities, builds the index, places facts where analysts look, and writes a ranked escalations file. |
| 3 | Query | Consult the knowledge wiki to answer deal-team questions — fresh session per question, read-only, every fact cited with provenance. |
| 4 | Operate | Periodically dream over the analyst session transcripts — the wiki reorganizes around real usage. Writes flow through the consolidation gate, never through query sessions. |
""")

md("""## Prerequisites

**Required:**
- Python 3.11+, and an Anthropic API key (set as `ANTHROPIC_API_KEY`) with access to the Claude Managed Agents beta
- **Dreaming research-preview access for your organization.** Dreaming is gated: if your org isn't enrolled, `/v1/dreams` returns 404 even with the beta header. Request access at [claude.com/form/claude-managed-agents](https://claude.com/form/claude-managed-agents) before you run Step 3.
- **The research-preview SDK build** (installed in the setup cell below) — the public PyPI `anthropic` package does not expose `client.beta.dreams` yet.
- Familiarity with the Managed Agents basics (agents, sessions, memory stores)

**Cost expectations, so nothing surprises you:** the default `mini` tier (26 documents, ~0.5 MB of text) cost about **\$35** to build end-to-end on the committed run — extraction sessions, one dream job, and the demo queries; treat that as an order-of-magnitude estimate, not a quote. Expect **roughly an hour of wall-clock**: extraction batches run concurrently, and the dream job took about 20 minutes on the committed run (larger corpora can take longer). Kick it off and check back. The `standard` tier (adds the merger proxy and tender documents) lands in the tens of dollars and reproduces study-scale results. Building once amortizes over every question you ask afterward: at full scale, wiki queries cost ~\$0.16 versus ~\$0.62 against raw documents, so the build pays for itself within a few hundred questions.
""")

code("""%%capture
%pip install -q httpx python-docx pillow python-dotenv""")

code('''# Dreaming ships in a dedicated research-preview SDK build; the public PyPI
# package does not include client.beta.dreams yet. Download the preview
# wheel from Anthropic's distribution link and verify its integrity before
# installing (the digest below is the build this notebook was tested with -
# if Anthropic publishes a new preview build, update both together).
import hashlib
import subprocess
import urllib.request

WHEEL_URL = "https://pkg.stainless.com/l/anthropic-python/5f5a2aac-6775-4d5f-bfac-fd747f5c661c"
WHEEL_SHA256 = "b92fe0480cd15f52830b572343e6e4b0be7a9c4eea058b64c9dfca958c4af539"
wheel_path = "/tmp/anthropic-0.100.0-py3-none-any.whl"

req = urllib.request.Request(WHEEL_URL, headers={"User-Agent": "kg-cookbook-setup/1.0"})
with urllib.request.urlopen(req) as resp, open(wheel_path, "wb") as f:
    f.write(resp.read())
digest = hashlib.sha256(open(wheel_path, "rb").read()).hexdigest()
assert digest == WHEEL_SHA256, f"wheel digest mismatch: {digest} - do not install; fetch an updated notebook"
subprocess.run(["python", "-m", "pip", "install", "-q", wheel_path], check=True)
print("research-preview SDK installed and verified")''')

code("""import base64
import json
import os
import time
from pathlib import Path

import httpx
from anthropic import Anthropic
from dotenv import load_dotenv

from utilities import poll_until_end_turn, wait_for_idle_status

load_dotenv()

# Per-stage model choice. The default (Sonnet 5 everywhere) is the
# quality-max configuration from the study behind this cookbook. In that
# study, a mixed configuration - a Fable-tier model for the build and an
# Opus-tier model for the dream - was about a third cheaper to build and
# ~12% cheaper per query at a one-point quality difference. Reproduce the
# method rather than the exact figures on your own corpus; sensible pairs
# to A/B against the default:
#   BUILD_MODEL = "claude-fable-5"     (builds the wiki)
#   DREAM_MODEL = "claude-opus-5"      (dreams / consolidates)
# Note: the dreaming API accepts sonnet- and opus-tier models only.
MODEL = os.environ.get("COOKBOOK_MODEL", "claude-sonnet-5")
BUILD_MODEL = os.environ.get("COOKBOOK_BUILD_MODEL", MODEL)   # wiki building
DREAM_MODEL = os.environ.get("COOKBOOK_DREAM_MODEL", MODEL)   # consolidation
QUERY_MODEL = os.environ.get("COOKBOOK_QUERY_MODEL", MODEL)   # analysts
DREAM_BETA = "dreaming-2026-04-21"   # research preview
API = "https://api.anthropic.com/v1"

# Long turns and busy hours are normal for this workload: generous request
# timeout, and aggressive retries so a transient 529 (overloaded) costs a
# backoff instead of your run.
client = Anthropic(timeout=httpx.Timeout(900.0, connect=10.0), max_retries=6)""")

md("""## Step 0a — Fetch the real document room

Three scripts (included, stdlib + pillow only) assemble the room from SEC EDGAR:

1. `build_manifest.py` — queries the EDGAR submissions index and writes tiered manifests (`mini` / `standard` / `full`)
2. `fetch_data_room.py --tier=mini` — downloads the filings and converts them to provenance-stamped text
3. `fetch_real_deck.py` — downloads six pages of Centerview Partners' board presentation (filed publicly as scanned slide images in the Schedule 13E-3 exhibits) and stitches them into a PDF

**In a hurry?** The `quickstart` tier is eight filings that still carry the whole deal arc — the original announcement, the revised terms that supersede it, the closing, and two insider filings — so entity resolution, temporal supersession, and consolidation all demonstrate in about 40 minutes end to end (validated: 41-minute wall-clock run). Every number in this notebook's saved outputs comes from the `mini` tier.

Every fetched document is a public filing; the header on each text file records its form type, filing date, accession number, and source URL. SEC asks for a descriptive `User-Agent` — set `EDGAR_USER_AGENT="your-name your-email"` first.
""")

code("""# Run once (network required). Uncomment to fetch here, or run in a terminal.
# !python3 build_manifest.py
# !python3 fetch_data_room.py --tier=mini        # 26 docs, the full walkthrough (~1h, ~$35)
# !python3 fetch_data_room.py --tier=quickstart  # 8 docs, the lunch-break path (~40 min, ~$25)
# !python3 fetch_real_deck.py

DOCS = Path("data_room/docs")
assert DOCS.exists() and any(DOCS.iterdir()), "run the fetch scripts first (see above)"
doc_files = sorted(DOCS.glob("*.txt"))
print(f"{len(doc_files)} real documents, {sum(f.stat().st_size for f in doc_files)/1e3:.0f} KB text")
[f.name for f in doc_files[:8]]""")

md("""## Step 0b — Normalize the PDF and Word documents

Real deal rooms are not `.txt`. The two format problems have different right answers:

**PDFs: send them to Claude natively.** The Messages API accepts PDF document blocks, and Claude sees both the text layer *and the rendered page images* — so a chart is read visually, values and all. Our board deck is the hard case on purpose: it's **scanned slides**, no text layer at all. We ask for a markdown transcript with two rules: tables become pipe-tables, and every figure becomes a bracketed description **with the values read off the chart**.

**Word documents: parse, don't rasterize.** `.docx` is a zip of XML — `python-docx` reads paragraphs and tables directly. (`make_analyst_docx.py` generated our sample: a marked-synthetic analyst note about this real deal, quoting EBITDA on the Credit Agreement basis — deliberately non-comparable to the company's reported figure.)

Both normalized documents join the corpus with the same provenance headers as the filings.

**Inventory your corpus by file type before you build.** The costliest ingestion failure is silent: a loader that only handles the formats you thought about quietly drops the rest, and the wiki looks complete while missing whole documents. (In our benchmark testing, one task's grading spec lived in an `.eml` file our first loader skipped - every downstream score suffered and nothing errored.) List the extensions in your data room first; spreadsheets (`openpyxl`), slide decks (`python-pptx`), and email (`email` stdlib) all parse in a few lines each, and anything exotic can go through the PDF path above.
""")

code('''NORMALIZE_SYSTEM = """You transcribe business documents into clean markdown
for a due-diligence corpus. Rules:
- Transcribe ALL substantive content faithfully; no summarizing, no commentary.
- Tables become markdown pipe-tables.
- EVERY chart or figure becomes a bracketed description that includes the
  values readable from it.
- Preserve any notes about metric definitions or bases verbatim - they matter.
- Keep document order."""


def normalize_pdf(path: Path) -> str:
    pdf_b64 = base64.standard_b64encode(path.read_bytes()).decode()
    resp = client.messages.create(
        model=BUILD_MODEL, max_tokens=8000, system=NORMALIZE_SYSTEM,
        messages=[{"role": "user", "content": [
            {"type": "document", "source": {"type": "base64", "media_type": "application/pdf", "data": pdf_b64}},
            {"type": "text", "text": "Transcribe this document to markdown per your rules."},
        ]}],
    )
    return "".join(b.text for b in resp.content if b.type == "text")


deck = Path("example_data/deal_room_real/real_board_deck.pdf")
deck_md = normalize_pdf(deck)
(DOCS / "board_deck_cii.txt").write_text(
    "[SOURCE: SEC EDGAR | form: SC 13E3 exhibit (c)(ii) | filed: 2024-06-17 | "
    "note: normalized from scanned slide images]\\n\\n" + deck_md)
print(deck_md[:1000])''')

md("""From scanned slides with no text layer, the transcript recovers the full deal-terms table — price per share, both termination fees, every financing commitment — and the offer-progression table. Those values existed only as ink on an image. (You'll also notice the deck calls the company "Pacific" and the sponsor "Phoenix": real board materials use code names, which is itself an entity-resolution lesson the wiki will handle.)
""")

code('''def normalize_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    parts = []
    for para in doc.paragraphs:
        if not para.text.strip():
            continue
        prefix = "# " if para.style.name.startswith("Heading 1") else ""
        parts.append(prefix + para.text)
    for table in doc.tables:
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        parts.append("| " + " | ".join(rows[0]) + " |")
        parts.append("|" + "---|" * len(rows[0]))
        for row in rows[1:]:
            parts.append("| " + " | ".join(row) + " |")
    return "\\n\\n".join(parts)


note_md = normalize_docx(Path("example_data/analyst_note.docx"))
(DOCS / "analyst_note.txt").write_text(
    "[SOURCE: synthetic demo document (marked) | analyst flash note | as-of: 2024-05-14]\\n\\n" + note_md)

# the vendor-feed CSV (also marked synthetic) rides along as-is - text formats are free
import shutil
shutil.copy("example_data/data_provider_extract.csv", DOCS / "data_provider_extract.csv.txt")

doc_files = sorted(DOCS.glob("*.txt"))
print(f"corpus complete: {len(doc_files)} documents")''')

md("""**Adapting to your own deal room** — route every file type explicitly so nothing is silently skipped:""")

code("""from pathlib import Path

def normalize_any(path: Path) -> str:
    \"\"\"One entry point per document; extend rather than filter.\"\"\"
    suffix = path.suffix.lower()
    if suffix in (".txt", ".md"):
        return path.read_text(errors="replace")
    if suffix == ".docx":
        return normalize_docx(path)          # python-docx: paragraphs + tables
    if suffix == ".pdf":
        return normalize_pdf(path)           # native PDF blocks: text + charts
    if suffix == ".xlsx":
        import openpyxl                      # pip install openpyxl
        wb = openpyxl.load_workbook(str(path), data_only=True)
        parts = []
        for ws in wb.worksheets:
            parts.append(f"## Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = ["" if v is None else str(v) for v in row]
                if any(c.strip() for c in cells):
                    parts.append(" | ".join(cells))
        return "\\n".join(parts)
    if suffix == ".pptx":
        from pptx import Presentation        # pip install python-pptx
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"## Slide {i}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
        return "\\n".join(parts)
    if suffix == ".eml":
        import email, email.policy           # stdlib
        msg = email.message_from_bytes(path.read_bytes(), policy=email.policy.default)
        hdr = "\\n".join(f"{k}: {msg.get(k, '')}" for k in ("From", "To", "Date", "Subject"))
        body = msg.get_body(preferencelist=("plain", "html"))
        return hdr + "\\n\\n" + (body.get_content() if body else "")
    # Anything else (scans, images, exotic formats): the PDF path reads pixels.
    raise ValueError(f"unhandled document type: {path.name} - add a branch, don't skip it")

# Example: corpus = {f.stem: normalize_any(f) for f in Path("my_deal_room").iterdir()}""")

md("""## Step 1 — Extraction: the schema rides on the store attachment

One agent; one session per **batch** of documents (small filings share a session; anything large is chunked). Two prompts do the work:

- The **system prompt** carries the rules our full-scale ablations showed are load-bearing: check for an existing entity before creating one (the deck's "Pacific" and the filings' "Squarespace" must land in one file), close superseded facts with `valid_to` instead of deleting them, and **never record a global absence claim** — only "not in this document". False absences were the dominant failure mode at scale, and they are nearly impossible to dislodge once filed.
- The **store attachment `instructions`** carry the wiki schema. Schema on the attachment, not the agent: the same agent can serve stores with different layouts.
""")

code('''environment = client.beta.environments.create(
    name="deal-kg-demo",
    config={"type": "cloud", "networking": {"type": "limited"}},
)

store = client.beta.memory_stores.create(
    name="Deal knowledge wiki",
    description="Knowledge wiki extracted from the take-private document room.",
)

EXTRACTOR_SYSTEM = """You are a deal room agent building a knowledge wiki for a
due-diligence team evaluating the take-private of Squarespace, Inc. (SQSP) by
Permira funds. You will be shown documents across many sessions. Your ONLY
output that matters is what you write into the attached memory store - the
wiki must be useful to a different agent later, with no access to this
conversation.

Rules:
- Each document begins with a [SOURCE: ...] provenance header - cite its
  form/slug in every fact you file.
- Before creating an entity file, check whether one already exists for the
  same real-world entity under a different name (list the entities/ dir);
  if so, extend it and add the alias instead of creating a duplicate.
  Board materials use code names (e.g. "Pacific" for the company) - resolve
  them to the canonical entity.
- Temporal discipline: when a later document supersedes an earlier fact
  (a changed price, a changed fee), do not delete the old fact - close it
  with valid_to and add the new one with valid_from.
- When a document amends, renews, or restates an earlier agreement, record
  each article the amendment TOUCHES section by section, every numbered
  section verbatim - including sections restated "for convenience" or
  carried over unchanged - and state what changed. Do not skip a
  carried-over section because your own reading says it is not implicated:
  file the provisions, not your synthesis of their net effect, so an
  analyst can derive interactions between new and old clauses that you did
  not anticipate. This full-article treatment applies ONLY to amended or
  restated articles; everywhere else the bounded-verbatim rule below
  governs.
- Quote verbatim ONLY operative language (consent standards, triggers,
  carve-outs, defined terms); summarize consequences and context. In our
  testing this bounded-verbatim rule cut build cost by a third at equal
  answer quality.
- Non-contract documents (filings, memos, disclosures) carry facts the
  contracts cannot: profile every affiliate, subsidiary, or portfolio
  company of any party (line of business, size) - cross-entity conflicts
  and carve-outs can only be assessed against these profiles.
- Never record a global absence claim; only "not in this document".
  Unanswered questions a deal team would ask go in open_questions/<doc-slug>.md.
- When you file a restrictive covenant (non-compete, exclusivity,
  non-solicit), add an open question on its enforceability under the stated
  governing law, naming the applicable legal standard (e.g. reasonableness
  review of restrictive covenants) - enforceability risk is invisible unless
  raised.
- Work silently and end with a one-line summary of what you added."""

KG_SCHEMA = """This store holds the deal knowledge wiki, in the Agent
Knowledge format: every node is a markdown file with YAML frontmatter, and
[[wikilinks]] in the body are the edges between nodes. Layout (paths
relative to the memory root):

entities/<slug>.md   one node per real-world entity. Frontmatter:
    ---
    name: <slug>            (matches the filename, the link target)
    description: <one line: what this entity is + why it matters here>
    aliases: [<every surface form seen, including code names>]
    metadata:
      type: company | person | fund | shell_entity | advisor | metric_definition | contract | event
    ---
    Body: bullet facts. EVERY fact ends with `[source: <doc> | as-of: <date>]`.
    When a fact mentions another entity in the store, write it as a
    [[wikilink]] to that entity's slug - links are the wiki's edges, and
    every link you add is a constraint that keeps neighbors honest.
    EVERY contract document (including each amendment, renewal, DPA, or
    side letter) owns a contract entity node whose slug is EXACTLY the
    document's slug from its [SOURCE] header - doc <slug> owns
    entities/<slug>.md, and ONLY the session reading that document writes
    that node. If a document you are reading merely mentions or summarizes
    a contract that is its own document elsewhere in the data room, do NOT
    create a stub node for it - reference it as a [[redlink]] and put the
    summary facts in your own document's files; the contract's own session
    fills the node. (Concurrent sessions overwrite each other's contract
    nodes wholesale otherwise - one summary stub silently destroyed a full
    verbatim node.) Record the contract's operative
    provisions VERBATIM in that node - never only as a pointer to an index,
    and never scattered across the party entities instead. Party entities
    summarize and [[link]] to the contract node. An analyst reading one
    node must not chase pointers for the exact language.

relations/<doc-slug>.md   relationships found in ONE source document,
    one per line, subjects and objects as [[wikilinks]]:
    `[[<subject-slug>]] | <relation> | [[<object-slug>]] | valid_from=<date> | valid_to=<date or open> | [source: <doc>]`

metrics/<doc-slug>.md   reported financial figures from ONE source document:
    `<metric> | <value> | basis/definition | period | [source: <doc>]`
    NEVER merge figures computed on different bases (e.g. two EBITDA
    definitions) - record each with its basis and keep both.

open_questions/<doc-slug>.md   gaps raised while reading that document.
    ONE line per question: name the gap - and, for enforceability
    questions, the governing law and the applicable legal standard - then
    stop. The question, not an essay; analysis belongs to the analyst who
    picks it up. A [[wikilink]] to an entity that does not exist yet is
    GOOD - a redlink marks the frontier for the resolver and the dream.
    Never guess at missing facts.

SINGLE-WRITER RULE (relations/, metrics/, open_questions/): write ONLY
under your own document's slugs - never a shared file or another
document's file. Extraction sessions run concurrently, and a shared
append-file invites one session to overwrite another's work wholesale
(we learned this the hard way: one wholesale rewrite silently destroyed
other sessions' verbatim provision text). Consolidation merges the
per-source files into unified views later.

ENTITY NODES ARE SHARED, keyed by real-world identity: one node per
entity, whichever sessions touch it. If a node for your contract or
party already exists - even a stub filed from a summary document -
extend THAT node: read-then-merge, never rewrite, and never a duplicate
node under a new slug."""

extractor = client.beta.agents.create(
    name="deal-kg-extractor",
    model=BUILD_MODEL,
    system=EXTRACTOR_SYSTEM,
    tools=[{
        "type": "agent_toolset_20260401",
        "default_config": {"enabled": True, "permission_policy": {"type": "always_allow"}},
    }],
)
print(extractor.id)''')

code('''# Pack documents into batches of ~120k characters (one session per batch);
# oversized documents get their own chunked sessions.
BATCH_CHARS = 120_000

batches, current, size = [], [], 0
for f in doc_files:
    text = f.read_text()
    if len(text) > BATCH_CHARS:                     # chunk the big ones
        for i in range(0, len(text), BATCH_CHARS):
            batches.append([(f"{f.stem}.part{i // BATCH_CHARS + 1}", text[i:i + BATCH_CHARS])])
        continue
    if size + len(text) > BATCH_CHARS and current:
        batches.append(current)
        current, size = [], 0
    current.append((f.stem, text))
    size += len(text)
if current:
    batches.append(current)

print(f"{len(doc_files)} documents -> {len(batches)} extraction sessions")''')

code('''# Launch every batch session up front - they run concurrently server-side -
# then wait for each. Wall-clock is roughly the slowest batch, not the sum.
# (Long turns are normal here: a session can even be rescheduled server-side
# mid-turn and still finish - the polling wait rides through it.)
extraction_sessions = []
for i, batch in enumerate(batches, 1):
    payload = "\\n\\n=====\\n\\n".join(text for _, text in batch)
    session = client.beta.sessions.create(
        agent={"type": "agent", "id": extractor.id, "version": extractor.version},
        environment_id=environment.id,
        resources=[{
            "type": "memory_store", "memory_store_id": store.id,
            "access": "read_write", "instructions": KG_SCHEMA,
        }],
    )
    client.beta.sessions.events.send(
        session_id=session.id,
        events=[{"type": "user.message",
                 "content": [{"type": "text", "text":
                     f"Read the following documents and file them into the wiki:\\n\\n{payload}"}]}],
    )
    print(f"launched batch {i}/{len(batches)}: {len(batch)} document(s) — "
          f"{', '.join(name for name, _ in batch[:3])}{', …' if len(batch) > 3 else ''}")
    extraction_sessions.append(session.id)

batch_payloads = {}
for sid, batch in zip(extraction_sessions, batches):
    batch_payloads[sid] = "\\n\\n=====\\n\\n".join(text for _, text in batch)

# Wait for every batch; if a turn ended abnormally (busy-hour server errors
# happen), re-send the same payload to the same session and wait again.
pending = list(extraction_sessions)
for attempt in range(1, 4):
    failed = []
    for i, sid in enumerate(pending, 1):
        print(f"--- waiting on batch session {i}/{len(pending)} (pass {attempt})")
        if poll_until_end_turn(client, sid) != "end_turn":
            failed.append(sid)
    if not failed:
        break
    print(f"re-sending {len(failed)} failed batch turn(s)")
    for sid in failed:
        client.beta.sessions.events.send(
            session_id=sid,
            events=[{"type": "user.message",
                     "content": [{"type": "text", "text":
                         "The previous attempt ended early. Read the following documents "
                         "and file them into the wiki (skip anything you already filed):"
                         f"\\n\\n{batch_payloads[sid]}"}]}],
        )
    pending = failed

len(extraction_sessions)''')

md("""## Step 2 — Resolve the wiki's open questions (in the build loop, not the query loop)

Extraction leaves an `open_questions.md`. At full scale, running a **resolver agent** on that list *before* the dream — rather than as a query-time fallback — was the single best architectural change we tested (96.2% vs 92.3%), because the dream then consolidates the researched facts exactly like extracted ones.

The resolver works in priority order: first cross-read the wiki itself (some questions are answerable from facts another session filed), then the public record. For genuinely unanswerable gaps it writes `CONFIRMED UNRESOLVABLE` and names the document a deal team should request — instead of guessing.
""")

code('''resolver_session = client.beta.sessions.create(
    agent={"type": "agent", "id": extractor.id, "version": extractor.version},
    environment_id=environment.id,
    resources=[{
        "type": "memory_store", "memory_store_id": store.id,
        "access": "read_write",
        "instructions": ("Read-write deal knowledge wiki. open_questions.md is your work "
                         "queue; entities/, relations.md, metrics.md hold the facts. Every "
                         "fact you add must carry [source | as-of] provenance."),
    }],
)
client.beta.sessions.events.send(
    session_id=resolver_session.id,
    events=[{"type": "user.message", "content": [{"type": "text", "text":
        "Work through open_questions.md (up to 5 highest-value questions). "
        "First try to RESOLVE each by cross-reading the wiki itself. Write "
        "resolved facts into the right entity/metrics file with provenance and "
        "mark the entry RESOLVED. For questions the corpus cannot answer, mark "
        "them CONFIRMED UNRESOLVABLE and name the document a deal team should "
        "request. Never guess."}]}],
)
poll_until_end_turn(client, resolver_session.id)''')

md("""## Step 3 — Reorganize the store via dreaming

Dreaming is offline consolidation: hand it the store plus the transcripts of every session that built it, and it writes a **new, reorganized store** — the input is untouched. That's not a detail, it's a discipline: every build is a version. Want to try different steering? Dream again, A/B the two stores on your own questions, keep the winner — instant rollback included. It deduplicates entities, builds an index, moves facts to where readers look, and recovers facts a session *read but never filed*.

The `instructions` string below was the highest-leverage prompt in our entire study — worth **+7.7 points** of answer quality by itself.

Two rules the ablations proved the hard way:

> **Dream the wiki as the sole source.** We tested telling the dream "analysts will also have the raw documents — write pointers where coverage is thin." It gutted the wiki: 21 points worse when queried alone. The dream heard *the reader has the documents* as license to not write things down, despite explicit instructions to keep coverage.
>
> **The dream never sees your corpus** — only the store and transcripts. It can reorganize what extraction learned; it cannot verify claims against source documents. Truth is won at extraction and resolution.
""")

code('''DREAM_RECORDS = []   # completed dream records, for the cost accounting


def run_dream(session_ids, store_id, instructions):
    """Start a dream over (sessions + store), watch it live, return the new store id.

    Uses the research-preview SDK's `client.beta.dreams` (installed in the
    setup cell). The dreaming API accepts a limited model list (sonnet and
    opus tiers). In the study behind this cookbook, an Opus-tier dream
    (claude-opus-5 on the current API) consolidated the same inputs at
    roughly half the dream cost of Sonnet (fewer passes) with slightly
    higher store quality - treat the ratio as directional and A/B on
    your own corpus. Live progress comes from two places: the dream's usage
    climbs as it writes, and the dream runs inside a session of its own -
    polling that shows consolidation threads spawning and files being
    written.
    """
    dream = client.beta.dreams.create(
        model=DREAM_MODEL,
        inputs=[
            {"type": "sessions", "session_ids": session_ids},
            {"type": "memory_store", "memory_store_id": store_id},
        ],
        instructions=instructions,
        betas=[DREAM_BETA],
    )
    print("dream:", dream.id)

    seen, threads, last_file = set(), 0, ""
    start = last_print = time.time()
    while True:
        d = client.beta.dreams.retrieve(dream.id, betas=[DREAM_BETA]).model_dump()
        if d["status"] not in {"queued", "in_progress", "running", "starting", "created", "pending"}:
            assert d["status"] == "completed", d
            DREAM_RECORDS.append(d)
            return d["outputs"][0]["memory_store_id"]
        try:
            for ev in client.beta.sessions.events.list(session_id=d["session_id"]):
                if ev.id in seen:
                    continue
                seen.add(ev.id)
                et = getattr(ev, "type", "")
                if et == "session.thread_created":
                    threads += 1
                elif et == "agent.tool_use":
                    path = str((getattr(ev, "input", None) or {}).get("file_path") or "")
                    if path:
                        last_file = path.split("/")[-1]
        except Exception:
            pass                      # progress peek is best-effort
        if time.time() - last_print >= 120:
            out_k = (d.get("usage") or {}).get("output_tokens", 0) / 1000
            bits = [f"dreaming… {(time.time() - start) / 60:.0f}m",
                    f"{out_k:.0f}k tokens written"]
            if threads:
                bits.append(f"{threads} consolidation thread(s)")
            if last_file:
                bits.append(f"working on {last_file}")
            print("  " + " · ".join(bits))
            last_print = time.time()
        time.sleep(30)


STEERING = ("In addition to your normal consolidation work: (1) PLACEMENT: every "
            "fact about a person or entity must appear in that entity's own file, even "
            "if that duplicates a fact recorded elsewhere - an analyst reading only one "
            "entity file must not miss facts about it. (2) Create an escalations.md file "
            "at the store root: a ranked list of every red-flag item a due-diligence team "
            "should escalate (unexplained ownership changes, insider trading activity "
            "during the pending deal, undisclosed legal-entity names, intra-document "
            "inconsistencies, metric-definition conflicts), each entry with a one-line "
            "why-it-matters and full provenance. (3) In each person's entity file, add a "
            "'## Red flags' section if any apply, citing sources."
    " (4) LINK STRUCTURE: the store uses the Agent Knowledge format - markdown "
    "nodes with YAML frontmatter, [[wikilinks]] as edges. Repair the link "
    "structure as you consolidate: every entity mention in a body should be a "
    "[[wikilink]]; a [[redlink]] to a missing node means create the node or "
    "record why it cannot exist; keep frontmatter name/aliases consistent with "
    "the filename."
)

kg_store_id = run_dream(extraction_sessions + [resolver_session.id], store.id, STEERING)
print("consolidated wiki:", kg_store_id)''')

code('''# What did the dream build? The index and the top of the ranked escalations.
listing = client.beta.memory_stores.memories.list(kg_store_id, view="full")
files = {m.path: m.content for m in listing}
print(sorted(files))
print()
print(files.get("/escalations.md", "")[:900])''')

md("""### The wiki is yours — export it

Everything the pipeline built is plain markdown in a memory store, and the API hands it back file by file. Export the consolidated wiki locally at any point — for inspection, versioning in your own repo, diffing two builds, or feeding downstream systems. (This is also your answer to "where does my data live": the store is queryable via the API, and a local export is one loop away.)
""")

code('''export_dir = Path("kg_export")
for m in listing:
    dest = export_dir / m.path.lstrip("/")
    dest.parent.mkdir(parents=True, exist_ok=True)
    dest.write_text(m.content)
print(f"exported {len(list(export_dir.rglob('*.md')))} wiki files to {export_dir}/")
sorted(str(f) for f in export_dir.rglob("*"))[:10]''')

md("""## Step 4 — Query the wiki, read-only

Analysts attach the consolidated store **read-only**, one fresh session per question. The attachment instructions route them through the index, require provenance on every fact, and script the miss behavior: say *"not in the data room"* and name the document needed — never guess.

What we deliberately do **not** do: give the analyst the raw documents as fallback. Our full 2×2 ablation measured that configuration — on the best wiki it cost 3.9 points of quality and 69% more tokens. Document access is *insurance against a bad wiki*, not an enhancement of a good one. Build the good wiki.

> **"But shouldn't the model write updates when it learns something?"** It should learn — through the build loop, not the query session. A query-time write bypasses every quality gate the wiki's facts passed (extraction rules, resolution, consolidation): one analyst's plausible inference becomes the next analyst's "fact." Instead: keep analysts read-only, optionally give them a tiny side store for *miss logging only* ("not in the data room" moments become candidate open questions), and periodically **dream over the analyst session transcripts** — they already carry everything analysts learned and struggled with. We measured this loop: a \$24 dream over eleven analyst transcripts halved subsequent hard-question read costs, and it records correct inferences with provenance through the normal consolidation gate. Analyst writes propose; the dream disposes.
""")

code('''import re

ANALYST_SYSTEM = """You are a due-diligence analyst agent. Your ONLY knowledge
source is the attached read-only memory store containing the deal knowledge
wiki. Answer questions by reading the wiki; cite the provenance the wiki
carries. If the wiki does not answer the question, say plainly "Not in the
data room" and name what document would be needed - never guess. Keep
answers tight: a few sentences plus citations."""

analyst = client.beta.agents.create(
    name="deal-analyst",
    model=QUERY_MODEL,
    system=ANALYST_SYSTEM,
    tools=[{
        "type": "agent_toolset_20260401",
        "default_config": {"enabled": True, "permission_policy": {"type": "always_allow"}},
    }],
)


from IPython.display import Markdown, display

analyst_sessions = []       # transcripts feed the usage dream in Step 5


def ask(question: str, wiki_store_id: str) -> dict:
    """One fresh read-only session per question. Returns usage stats."""
    session = client.beta.sessions.create(
        agent={"type": "agent", "id": analyst.id, "version": analyst.version},
        environment_id=environment.id,
        resources=[{
            "type": "memory_store", "memory_store_id": wiki_store_id,
            "access": "read_only",
            "instructions": ("Read-only deal knowledge wiki. Start at _index.md; "
                             "escalations.md ranks red flags; facts carry "
                             "[source | as-of] provenance - cite it. If the wiki does "
                             "not answer, say 'Not in the data room' and name the "
                             "document needed."),
        }],
    )
    client.beta.sessions.events.send(
        session_id=session.id,
        events=[{"type": "user.message", "content": [{"type": "text", "text": question}]}],
    )
    reply: list = []
    stats: dict = {}
    poll_until_end_turn(client, session.id, sink=reply, stats=stats)
    analyst_sessions.append(session.id)
    answer = "\\n\\n".join(reply).replace("$", "\\\\$")   # keep MathJax away from prices
    # "- #1 ..." after a bullet renders as a giant heading; escape the hash
    answer = re.sub(r"(?m)^(\\s*[-*]\\s+)#", r"\\1\\\\#", answer)
    display(Markdown(f"**Q: {question}**\\n\\n**🧠 Analyst agent's answer** *(generated from the wiki — everything until the rule below)*:\\n\\n{answer}\\n\\n---"))
    print(f"  [{stats.get('tokens_in', 0):,} tokens in · {stats.get('reads', 0)} file reads]")
    return stats''')

code('''# Temporal supersession across REAL filings: the answer must lead with \$46.50
# and know the path from \$44.00 (announce 8-K -> amendment 8-K).
q_price = ask("What is the per-share merger consideration, and how did it change during the process?", kg_store_id)''')

code('''# Cross-format metric reconciliation: the company's reported figure (filings)
# vs the Credit Agreement estimate (Word analyst note) - non-comparable bases.
q_ebitda = ask("I've seen more than one FY2023 adjusted EBITDA figure in this room. Reconcile them - are they comparable?", kg_store_id)''')

code('''# Founder entity resolution: one person, multiple vehicles, plus supervoting.
q_control = ask("Who controls the company's voting power, and through which entities?", kg_store_id)''')

code('''# The trap: post-closing financials are NOT in the room (the company was
# delisted). A correct answer refuses and names what would be needed.
q_trap = ask("What was the company's revenue for fiscal year 2025?", kg_store_id)''')

md("""## Step 5 — Dream over the analyst transcripts: the wiki learns from being used

The write-back loop, closed safely. The four query sessions above are transcripts now — they carry what analysts looked for, where they had to hop between files, and what they inferred. A second dream over **(analyst transcripts + current wiki)** consolidates that experience: it reorganizes so the questions analysts actually ask are answerable more directly, and it can record sound inferences with provenance — all through the same quality gate as every other fact.

At study scale this was finding 11: a \$24 usage dream left answer *quality* unchanged (it cannot add facts nobody filed) but **halved the reading cost** of hard questions. Usage dreams buy organization, and the savings compound with every question asked afterward.
""")

code('''USAGE_STEERING = ("In addition to your normal consolidation work: the input sessions "
                  "are analyst Q&A transcripts run against this wiki. Study where "
                  "analysts struggled - searched several files, hedged, or synthesized "
                  "an answer from scattered fragments - and reorganize the wiki so "
                  "those and similar questions are answerable directly next time "
                  "(placement, indexing, cross-references between entity files). Where "
                  "an analyst derived a correct fact by sound inference over the wiki, "
                  "record it explicitly with provenance to the underlying sources. Keep "
                  "the escalations.md ranking current.")

kg2_store_id = run_dream(analyst_sessions, kg_store_id, USAGE_STEERING)
print("post-usage wiki:", kg2_store_id)''')

md("""### Round 2 — same question, improved wiki

Re-ask the price-evolution question against the post-usage store, then compare the cost of answering. One run is one sample — your numbers will wobble — but the direction to expect (and what we measured at scale) is the same answer for fewer file reads and fewer tokens, because the dream moved the answer to where the analyst looks first.
""")

code('''q_price_2 = ask("What is the per-share merger consideration, and how did it change during the process?", kg2_store_id)''')

code('''display(Markdown(f"""
| price-evolution question | tokens in | file reads |
|---|---|---|
| round 1 (post-build wiki) | {q_price["tokens_in"]:,} | {q_price["reads"]} |
| round 2 (post-usage wiki) | {q_price_2["tokens_in"]:,} | {q_price_2["reads"]} |

*(n=1 per cell — direction matters more than magnitude; at study scale the
usage dream halved hard-question reading cost across the battery.)*"""))''')

code('''# And a question the first round never asked - the usage dream should not
# have hurt anything it did not touch.
q_fees = ask("Summarize the termination fee structure and how it changed with the amendment.", kg2_store_id)''')

md("""## The economics of this run, measured

Everything above emitted usage events, so the notebook can price itself — no estimates for anything we ran. The one thing worth measuring that we *haven't* run yet is the counterfactual: what does a question cost **without** the wiki? Rather than quote our study, measure it here: load the raw corpus into a store and ask the same price-evolution question against the documents, no graph.

(The price sheet below is the only hardcoded input: public list rates per MTok for the default model. Swap in your own numbers if you run a different model or negotiated pricing.)
""")

code('''# one raw-documents store: each corpus file uploaded as-is (chunk the big ones)
raw_store = client.beta.memory_stores.create(
    name="Raw documents (baseline)",
    description="The un-extracted document room, for the no-wiki cost baseline.")

MAX_CHARS = 85_000
for f in doc_files:
    text = f.read_text()
    if len(text) <= MAX_CHARS:
        client.beta.memory_stores.memories.create(raw_store.id, path=f"/{f.stem}.txt", content=text)
    else:
        for i in range(0, len(text), MAX_CHARS):
            client.beta.memory_stores.memories.create(
                raw_store.id, path=f"/{f.stem}.part{i // MAX_CHARS + 1}.txt", content=text[i:i + MAX_CHARS])

raw_session = client.beta.sessions.create(
    agent={"type": "agent", "id": analyst.id, "version": analyst.version},
    environment_id=environment.id,
    resources=[{
        "type": "memory_store", "memory_store_id": raw_store.id,
        "access": "read_only",
        "instructions": ("Raw deal data-room documents as text files (large filings are "
                         "split into .partN files). Search and read them to answer "
                         "questions; cite the source files. If the documents do not "
                         "answer, say 'Not in the data room' - never guess."),
    }],
)
client.beta.sessions.events.send(
    session_id=raw_session.id,
    events=[{"type": "user.message", "content": [{"type": "text", "text":
        "What is the per-share merger consideration, and how did it change during the process?"}]}],
)
raw_reply, raw_stats = [], {}
poll_until_end_turn(client, raw_session.id, sink=raw_reply, stats=raw_stats)
print(f"raw-documents answer costs: {raw_stats.get('tokens_in', 0):,} tokens in · {raw_stats.get('reads', 0)} file reads")''')

code('''# Price this run. Only PRICES is hardcoded - public list rates per MTok.
PRICES = {"input": 2.00, "cache_write": 2.50, "cache_read": 0.20, "output": 10.00}


def usage_usd(u: dict) -> float:
    return (u.get("input", 0) * PRICES["input"]
            + u.get("cache_write", 0) * PRICES["cache_write"]
            + u.get("cache_read", 0) * PRICES["cache_read"]
            + u.get("output", 0) * PRICES["output"]) / 1e6


def session_usage(session_id: str) -> dict:
    u = {"input": 0, "cache_write": 0, "cache_read": 0, "output": 0}
    for ev in client.beta.sessions.events.list(session_id=session_id):
        if getattr(ev, "type", "") == "span.model_request_end":
            mu = getattr(ev, "model_usage", None)
            if mu is not None:
                u["input"] += getattr(mu, "input_tokens", 0) or 0
                u["cache_write"] += getattr(mu, "cache_creation_input_tokens", 0) or 0
                u["cache_read"] += getattr(mu, "cache_read_input_tokens", 0) or 0
                u["output"] += getattr(mu, "output_tokens", 0) or 0
    return u


def dream_usage(record: dict) -> dict:
    du = record.get("usage") or {}
    return {"input": du.get("input_tokens", 0), "cache_write": du.get("cache_creation_input_tokens", 0),
            "cache_read": du.get("cache_read_input_tokens", 0), "output": du.get("output_tokens", 0)}


extract_usd = sum(usage_usd(session_usage(s)) for s in extraction_sessions)
resolve_usd = usage_usd(session_usage(resolver_session.id))
dream1_usd = usage_usd(dream_usage(DREAM_RECORDS[0]))
dream2_usd = usage_usd(dream_usage(DREAM_RECORDS[1])) if len(DREAM_RECORDS) > 1 else 0.0
build_usd = extract_usd + resolve_usd + dream1_usd

wiki_queries = [q_price, q_ebitda, q_control, q_trap, q_price_2, q_fees]
wiki_q_usd = sum(usage_usd(q) for q in wiki_queries) / len(wiki_queries)
raw_q_usd = usage_usd(raw_stats)
breakeven = build_usd / max(raw_q_usd - wiki_q_usd, 1e-9)

display(Markdown(f"""
| what | cost (this run, measured) |
|---|---|
| build the wiki (extraction ${extract_usd:.2f} + resolver ${resolve_usd:.2f} + dreaming ${dream1_usd + dream2_usd:.2f}) | **${build_usd + dream2_usd:.2f}** one-time |
| ask a question **with** the wiki (mean of {len(wiki_queries)}) | **${wiki_q_usd:.3f}** |
| ask the same question against **raw documents** (baseline) | **${raw_q_usd:.3f}** |
| build breaks even after | **~{breakeven:.0f} questions** |"""))''')

md("""## What this buys you at scale

The mini tier proves the mechanics on real documents. At full study scale (the same transaction, ~2.7 MB corpus, a 21-question judged battery, 13 system variants):

| configuration | answer quality | tokens/question | cost/question |
|---|---|---|---|
| **this pipeline** (extract → resolve → steered dream, queried wiki-only) | **96.2%** | **68.5k** | **\$0.16** |
| agentic search over the raw documents | 85.7% | 284.5k | \$0.62 |

**Design rules distilled from the ablations:**

1. **Truth is won at extraction.** A fact never filed is rarely recovered downstream. Spend your rules there.
2. **Resolve gaps before the dream, not after.** Researched facts get indexed and placed like native ones.
3. **Steer the dream.** One instructions string was the biggest single quality lever.
4. **Dream sole-source; query sole-source.** Telling the dream about fallback documents gutted the wiki; giving analysts fallback documents taxed every question. The document budget belongs in the build loop.
5. **Treat absence claims as hypotheses.** Every layer — extraction, query, judging, even our own answer keys — made false "not disclosed" claims at some point. Verify absences mechanically.

**Scaling up from here:** rerun `fetch_data_room.py --tier=standard` for the 37-document room (adds the merger proxy and tender documents — and with them, answers to questions the mini room correctly refuses). The pipeline is unchanged; only the corpus grows.
""")

md("""## Clean up

Delete the demo resources. (In production the consolidated wiki store is the artifact you keep.)
""")

code('''for sid in extraction_sessions + [resolver_session.id] + analyst_sessions + [raw_session.id]:
    wait_for_idle_status(client, sid)
    client.beta.sessions.archive(sid)

client.beta.memory_stores.delete(store.id)
client.beta.memory_stores.delete(kg_store_id)
client.beta.memory_stores.delete(kg2_store_id)
client.beta.memory_stores.delete(raw_store.id)
client.beta.agents.archive(extractor.id)
client.beta.agents.archive(analyst.id)
client.beta.environments.archive(environment.id)
print("cleaned up")''')

md("""## Recap

We turned a real, mixed-format document corpus into a queryable knowledge wiki — using an M&A data room as the worked example, with a pipeline that transfers to any of the scenarios in the introduction:

- **Fetched** 26 public filings from EDGAR, reproducibly, at a chosen cost tier
- **Normalized** a scanned board deck (values recovered from chart images) and a Word analyst note into the corpus
- **Extracted** into a memory store with the schema on the attachment and the anti-absence rules in the system prompt
- **Resolved** open questions in the build loop, honestly marking the unresolvable
- **Consolidated** with one steered, sole-source dream
- **Queried** read-only with provenance — surviving a real price supersession, a metric-basis conflict, a code-name entity resolution, and an absence trap
- **Closed the loop**: dreamed over the analyst transcripts so the wiki reorganizes around real usage — writes flow through the consolidation gate, never through query sessions

Two field lessons worth carrying: for report-scale deliverables (memos, full analyses), tell the agent *"your reply is the deliverable"* — agents with file tools will otherwise write the report to a workspace file and hand you a summary. And when you evaluate variants (steering tweaks, prompt changes), replicate before believing: identical configurations spread by several rubric points run-to-run — a single-run delta is noise until it repeats.

This notebook is a preview we're sharing for feedback: what worked, what confused you, what your own document rooms need that this doesn't cover. Tell us.
""")

nb = {
    "cells": cells,
    "metadata": {
        "kernelspec": {"display_name": "Python 3", "language": "python", "name": "python3"},
        "language_info": {"name": "python", "version": "3.11.0"},
    },
    "nbformat": 4,
    "nbformat_minor": 5,
}
out = Path(__file__).parent / "distill_documents_into_knowledge_graph.ipynb"
json.dump(nb, open(out, "w"), indent=1)
print(f"wrote {out}: {len(cells)} cells")
